import json
import os
import requests
import tempfile
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from io import BytesIO
import logging
from utils.cloudinary import upload_to_cloudinary, store_presentation_url
from utils.openai import openai_client
from config import settings
logger = logging.getLogger(__name__)    

class SlideContent:
    def __init__(self, title: str, content: List[str], slide_type: str = "content", image_query: str = None, layout: str = "content"):
        self.title = title
        self.content = content
        self.slide_type = slide_type
        self.image_query = image_query or title     
        self.layout = layout

def download_image_from_pexels(query: str, width: int = 800, height: int = 600) -> BytesIO:
    try:
        api_key = settings.PEXELS_API_KEY
        print("API Key:", api_key)
        print("Query:", query)
        url = f"https://api.pexels.com/v1/search?query={query.replace(' ', '%20')}&per_page=1"
        headers = {"Authorization": api_key}
        
        response = requests.get(url, headers=headers, timeout=10)
        if response.status_code == 200:
            photos = response.json().get("photos", [])
            if photos:
                image_url = photos[0]["src"]["medium"]
                image_response = requests.get(image_url, timeout=10)
                if image_response.status_code == 200:
                    return BytesIO(image_response.content)
        logger.warning(f"Failed to download image for query: {query} (Status: {response.status_code})")
        return None
    except Exception as e:
        logger.error(f"Error downloading image from Pexels: {e}")
        return None

def download_image_from_pixabay(query: str, width: int = 800, height: int = 600) -> BytesIO:
    try:
        api_key = settings.PIXABAY_API_KEY
        url = f"https://pixabay.com/api/?key={api_key}&q={query.replace(' ', '%20')}&image_type=photo&per_page=3"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            hits = response.json().get("hits", [])
            if hits:
                image_url = hits[0]["webformatURL"]
                image_response = requests.get(image_url, timeout=10)
                if image_response.status_code == 200:
                    return BytesIO(image_response.content)
        logger.warning(f"Failed to download image for query: {query} (Status: {response.status_code})")
        return None
    except Exception as e:
        logger.error(f"Error downloading image from Pixabay: {e}")
        return None

def download_image_from_unsplash(query: str, width: int = 800, height: int = 600) -> BytesIO:
    try:
        url = f"https://source.unsplash.com/{width}x{height}/?{query.replace(' ', '%20')}"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            return BytesIO(response.content)
        else:
            logger.warning(f"Failed to download image for query: {query}")
            return None
    except Exception as e:
        logger.error(f"Error downloading image: {e}")
        return None

def download_image(query: str, width: int = 800, height: int = 600) -> BytesIO:
    image_stream = download_image_from_pexels(query, width, height)
    if image_stream:
        return image_stream
    
    image_stream = download_image_from_pixabay(query, width, height)
    if image_stream:
        return image_stream
    
    return download_image_from_unsplash(query, width, height)

def style_text_box(text_frame, font_size=18, is_title=False):
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = 'Calibri'
        paragraph.font.size = Pt(font_size)
        
        if is_title:
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
        else:
            paragraph.font.color.rgb = RGBColor(64, 64, 64)
            paragraph.space_after = Pt(6)

def create_title_slide(prs, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_right = Inches(0.5)
    
    for paragraph in title_frame.paragraphs:
        paragraph.font.name = 'Calibri'
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(31, 78, 121)
        paragraph.alignment = PP_ALIGN.CENTER
    
    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_frame = subtitle_shape.text_frame
        
        for paragraph in subtitle_frame.paragraphs:
            paragraph.font.name = 'Calibri'
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = RGBColor(64, 64, 64)
            paragraph.alignment = PP_ALIGN.CENTER
    
    return slide

def create_agenda_slide(prs, slide_data: SlideContent):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(9)
    title_height = Inches(1)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = 'Calibri'
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(31, 78, 121)
    
    content_left = Inches(1.5)
    content_top = Inches(2)
    content_width = Inches(7)
    content_height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.clear()
    
    for i, point in enumerate(slide_data.content):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = f"{i+1}. {point}"
        p.font.name = 'Calibri'
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(64, 64, 64)
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

def create_content_slide_with_image(prs, slide_data: SlideContent):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    image_stream = download_image(slide_data.image_query)
    
    if image_stream:
        try:
            img_left = Inches(5.5)
            img_top = Inches(1.5)
            img_width = Inches(4)
            img_height = Inches(3)
            
            slide.shapes.add_picture(image_stream, img_left, img_top, img_width, img_height)
        except Exception as e:
            logger.warning(f"Could not add image: {e}")
    
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(9)
    title_height = Inches(1)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = 'Calibri'
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(31, 78, 121)
    
    content_left = Inches(0.5)
    content_top = Inches(2)
    content_width = Inches(4.5)
    content_height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.clear()
    
    for i, point in enumerate(slide_data.content):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = f"• {point}"
        p.font.name = 'Calibri'
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(64, 64, 64)
        p.space_after = Pt(8)
        p.level = 0
    
    return slide

def create_two_column_slide(prs, slide_data: SlideContent):
    """Create a two-column content slide"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.8)
    title_width = Inches(9)
    title_height = Inches(1)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = slide_data.title
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = 'Calibri'
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(31, 78, 121)
    
    # Split content into two columns
    mid_point = len(slide_data.content) // 2
    left_content = slide_data.content[:mid_point]
    right_content = slide_data.content[mid_point:]
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(4.5))
    left_frame = left_box.text_frame
    left_frame.clear()
    
    for i, point in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.name = 'Calibri'
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(64, 64, 64)
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4.5), Inches(4.5))
    right_frame = right_box.text_frame
    right_frame.clear()
    
    for i, point in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.name = 'Calibri'
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(64, 64, 64)
        p.space_after = Pt(8)
    
    return slide

def create_section_slide(prs, title: str):
    """Create a section divider slide"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add centered title
    title_left = Inches(1)
    title_top = Inches(3)
    title_width = Inches(8)
    title_height = Inches(2)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = 'Calibri'
    title_paragraph.font.size = Pt(48)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(31, 78, 121)
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    return slide

def create_powerpoint(slides: List[SlideContent], presentation_id: str, topic: str) -> str:
    """Create PowerPoint presentation from slides with enhanced styling and images"""
    try:
        prs = Presentation()
        
        # Remove default slide
        if len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        # Process slides
        section_count = 0
        
        for i, slide_data in enumerate(slides):
            
            if slide_data.slide_type == "title":
                # Create title slide
                create_title_slide(prs, slide_data.title, "Professional Presentation")
                
            elif slide_data.slide_type == "agenda":
                # Create agenda slide
                create_agenda_slide(prs, slide_data)
                
            elif slide_data.slide_type == "section":
                # Create section divider
                create_section_slide(prs, slide_data.title)
                section_count += 1
                
            elif slide_data.layout == "two_column":
                # Create two-column slide
                create_two_column_slide(prs, slide_data)
                
            else:
                # Create content slide with image
                create_content_slide_with_image(prs, slide_data)
                
                # Add section slides every 3-4 content slides for longer presentations
                if (i > 0 and i % 5 == 0 and section_count < 2 and len(slides) > 10):
                    section_title = f"Section {section_count + 1}"
                    create_section_slide(prs, section_title)
                    section_count += 1
        
        # Add a thank you slide at the end
        thank_you_slide = SlideContent(
            title="Thank You",
            content=["Questions & Discussion", "Contact for more information"],
            slide_type="content",
            image_query="thank you business meeting"
        )
        create_content_slide_with_image(prs, thank_you_slide)
        
        # Save presentation to a temporary file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            prs.save(temp_file.name)
            
            # Upload to Cloudinary
            cloudinary_url = upload_to_cloudinary(temp_file.name, presentation_id, topic)
            
            if cloudinary_url:
                # Store URL in Redis
                store_presentation_url(presentation_id, cloudinary_url)
                
                # Clean up temporary file
                os.unlink(temp_file.name)
                
                return cloudinary_url
            else:
                raise Exception("Failed to upload presentation to Cloudinary")
        
    except Exception as e:
        logger.error(f"Error creating PowerPoint: {e}")
        raise Exception(f"Failed to create presentation: {str(e)}")

def _get_presentation_structure(slide_count: int, presentation_type: str) -> str:
    """Define presentation structure based on slide count and type"""
    
    if slide_count <= 5:
        return "1. Title slide, 2. Problem/Challenge, 3. Solution/Approach, 4. Benefits/Results, 5. Next Steps"
    elif slide_count <= 10:
        return "1. Title slide, 2. Agenda, 3. Background/Context, 4-7. Main Content (key points), 8. Recommendations, 9. Implementation, 10. Q&A"
    elif slide_count <= 15:
        return "1. Title slide, 2. Agenda, 3. Executive Summary, 4. Background, 5-11. Main Content (detailed analysis), 12. Recommendations, 13. Implementation Plan, 14. Next Steps, 15. Q&A"
    else:
        return "1. Title slide, 2. Agenda, 3. Executive Summary, 4-6. Background & Context, 7-15. Detailed Analysis (multiple sections), 16-18. Recommendations & Strategy, 19. Implementation Roadmap, 20. Q&A"

def _generate_realistic_fallback_slides(topic: str, slide_count: int) -> List[SlideContent]:
    """Generate realistic fallback slides when AI generation fails"""
    
    slides = []
    
    # Title slide
    slides.append(SlideContent(
        title=f"{topic}: Strategic Overview",
        content=[],
        slide_type="title",
        image_query=f"{topic} professional business",
    ))
    
    # Agenda slide (if enough slides)
    if slide_count > 5:
        slides.append(SlideContent(
            title="Agenda",
            content=[
                "Current market landscape analysis",
                "Key challenges and opportunities", 
                "Strategic recommendations",
                "Implementation roadmap",
                "Success metrics and KPIs",
                "Q&A and discussion"
            ],
            slide_type="agenda",
            image_query="business agenda meeting",
        ))
    
    # Generate main content slides with realistic topics
    content_slides = slide_count - 3 if slide_count > 5 else slide_count - 2
    
    content_topics = [
        ("Market Analysis & Trends", "market analysis data charts", "Current market dynamics and emerging trends"),
        ("Key Challenges", "business challenges obstacles", "Primary obstacles and pain points"),
        ("Strategic Opportunities", "business opportunities growth", "Growth opportunities and market gaps"),
        ("Recommended Approach", "strategy planning roadmap", "Our strategic approach and methodology"),
        ("Implementation Framework", "project management timeline", "Step-by-step implementation plan"),
        ("Resource Requirements", "team resources budget allocation", "Required resources and investments"),
        ("Risk Assessment", "risk management analysis", "Potential risks and mitigation strategies"),
        ("Success Metrics", "KPI dashboard analytics", "Key performance indicators and targets"),
        ("Expected Outcomes", "business results success", "Projected results and benefits"),
        ("Timeline & Milestones", "project timeline milestones", "Implementation schedule and key dates")
    ]
    
    for i in range(min(content_slides, len(content_topics))):
        title_suffix, image_query, description = content_topics[i]
        
        # Create more realistic content points
        content_points = [
            f"Critical success factors for {topic.lower()}",
            f"Data-driven insights and analysis",
            f"Best practices from industry leaders", 
            f"Quantifiable impact on business metrics",
            f"Alignment with organizational objectives"
        ]
        
        # Add specific content based on slide topic
        if "Market" in title_suffix:
            content_points = [
                f"Market size: $X billion growing at Y% CAGR",
                f"Key competitors and market positioning",
                f"Customer segments and demographics",
                f"Emerging trends and disruptions",
                f"Regulatory environment and compliance"
            ]
        elif "Challenge" in title_suffix:
            content_points = [
                f"Primary operational inefficiencies",
                f"Resource constraints and limitations",
                f"Technology gaps and legacy systems",
                f"Competitive pressure and market dynamics",
                f"Regulatory and compliance requirements"
            ]
        elif "Opportunit" in title_suffix:
            content_points = [
                f"Untapped market segments worth $X million",
                f"Process optimization potential: 20-30% efficiency gains",
                f"Technology adoption opportunities",
                f"Strategic partnerships and alliances",
                f"Innovation and product development areas"
            ]
        
        slide_layout = "two_column" if len(content_points) > 5 else "content"
        
        slides.append(SlideContent(
            title=f"{title_suffix}: {topic}",
            content=content_points,
            slide_type="content",
            image_query=image_query,
            layout=slide_layout
        ))
    
    # Conclusion slide
    slides.append(SlideContent(
        title="Key Takeaways & Next Steps",
        content=[
            f"{topic} presents significant growth opportunities",
            "Strategic approach will deliver measurable ROI",
            "Implementation requires cross-functional collaboration",
            "Success depends on consistent execution",
            "Continuous monitoring and optimization essential"
        ],
        slide_type="content",
        image_query=f"{topic} success team collaboration",
    ))
    
    return slides[:slide_count]

def generate_presentation_content(topic: str, slide_count: int, presentation_type: str = "business"):
    """Generate enhanced content with better structure"""
    print("Generating enhanced content for topic:", topic)
    print("Slide count:", slide_count)
    print("Presentation type:", presentation_type)
    
    try:
        # Get structure based on slide count
        structure = _get_presentation_structure(slide_count, presentation_type)
        
        prompt = f"""
        Create a professional {presentation_type} presentation with {slide_count} slides about: {topic}
        
        Follow this structure: {structure}
        
        Guidelines for realistic PowerPoint content:
        - Use concise, impactful bullet points (5-8 words max per point)
        - Include relevant statistics, facts, or data points where appropriate
        - Add actionable insights and recommendations
        - Use professional business language
        - Create compelling slide titles (6-10 words)
        - Suggest appropriate slide layouts
        
        For each slide, provide:
        - Compelling title that captures the key message
        - 3-6 bullet points with substantial, realistic content
        - Appropriate slide layout (content, two_column, agenda, etc.)
        - Relevant image search query
        
        Return in JSON format:
        {{
            "presentation_title": "Professional Title",
            "subtitle": "Engaging subtitle",
            "slides": [
                {{
                    "title": "Slide Title",
                    "content": ["Concise bullet point 1", "Impactful bullet point 2"],
                    "slide_type": "title|content|section|agenda",
                    "layout": "content|two_column|agenda",
                    "image_query": "professional search terms",
                }}
            ]
        }}
        """
        
        response = openai_client.chat.completions.create(
            model="gpt-4",  # Using GPT-4 for better quality
            messages=[
                {
                    "role": "system", 
                    "content": "You are an expert presentation designer with 10+ years creating executive-level PowerPoint presentations. Focus on clarity, impact, and professional appeal."
                },
                {"role": "user", "content": prompt}
            ],
            max_tokens=4000,
            temperature=0.6
        )
        
        content = response.choices[0].message.content
        
        try:
            # Parse the JSON response
            slides_data = json.loads(content)
            slides = []
            
            for slide_info in slides_data.get("slides", []):
                slide = SlideContent(
                    title=slide_info.get("title", "Untitled Slide"),
                    content=slide_info.get("content", ["Content not available"]),
                    slide_type=slide_info.get("slide_type", "content"),
                    image_query=slide_info.get("image_query", f"{topic} professional"),
                    layout=slide_info.get("layout", "content")
                )
                slides.append(slide)
            
            return slides
            
        except json.JSONDecodeError as e:
            logger.warning(f"JSON parsing failed: {e}. Using realistic fallback content.")
            return _generate_realistic_fallback_slides(topic, slide_count)
            
    except Exception as e:
        logger.error(f"Error generating enhanced content: {e}")
        return _generate_realistic_fallback_slides(topic, slide_count)