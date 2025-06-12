import os
import requests
import tempfile
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from io import BytesIO
import logging
from .storage import upload_to_cloudinary, store_presentation_url
from .helper import create_presentation_content

logger = logging.getLogger(__name__)

def add_gradient_background(slide, start_color, end_color):
    """Add a professional gradient background to the slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(*start_color)
    fill.gradient_stops[1].color.rgb = RGBColor(*end_color)

def download_image_from_unsplash(query: str) -> Optional[BytesIO]:
    """Download a relevant image from Unsplash"""
    try:
        response = requests.get(
            f"https://source.unsplash.com/featured/?{query}",
            timeout=10
        )
        if response.status_code == 200:
            return BytesIO(response.content)
    except Exception as e:
        logger.warning(f"Could not download image: {e}")
    return None

def create_title_slide(prs, title: str, subtitle: str = ""):
    """Create an attractive title slide with modern styling"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Add gradient background
    add_gradient_background(slide, (31, 78, 121), (100, 149, 237))
    
    # Style title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.margin_left = Inches(0.5)
    title_frame.margin_right = Inches(0.5)
    
    for paragraph in title_frame.paragraphs:
        paragraph.font.name = 'Calibri'
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Add subtitle if provided
    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_frame = subtitle_shape.text_frame
        
        for paragraph in subtitle_frame.paragraphs:
            paragraph.font.name = 'Calibri'
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide_with_image(prs, slide_data: Dict):
    """Create content slide with image and professional layout"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add subtle gradient background
    add_gradient_background(slide, (245, 245, 245), (255, 255, 255))
    
    # Download and add image
    image_stream = download_image_from_unsplash(slide_data.get('image_query', slide_data['title']))
    
    if image_stream:
        try:
            # Add image on the right side
            img_left = Inches(5.5)
            img_top = Inches(1.5)
            img_width = Inches(4)
            img_height = Inches(3)
            
            slide.shapes.add_picture(image_stream, img_left, img_top, img_width, img_height)
        except Exception as e:
            logger.warning(f"Could not add image: {e}")
    
    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(9)
    title_height = Inches(1)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = slide_data['title']
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = 'Calibri'
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(31, 78, 121)
    
    # Add content
    content_left = Inches(0.5)
    content_top = Inches(2)
    content_width = Inches(4.5) if image_stream else Inches(9)
    content_height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Clear default paragraph
    content_frame.clear()
    
    for i, point in enumerate(slide_data['content']):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = f"• {point}"
        p.font.name = 'Calibri'
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(64, 64, 64)
        p.space_after = Pt(8)
        p.level = 0
    
    return slide

def create_section_slide(prs, title: str):
    """Create a section divider slide with modern design"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add gradient background
    add_gradient_background(slide, (70, 130, 180), (135, 206, 235))
    
    # Add centered title
    title_left = Inches(1)
    title_top = Inches(3)
    title_width = Inches(8)
    title_height = Inches(2)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.name = 'Calibri'
    title_paragraph.font.size = Pt(48)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(255, 255, 255)
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    return slide

def create_powerpoint(slides: List[Dict], presentation_id: str) -> str:
    """Create PowerPoint presentation from slides with enhanced styling and images"""
    try:
        prs = Presentation()
        
        # Remove default slide
        if len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        
        # Process slides
        section_count = 0
        
        for i, slide_data in enumerate(slides):
            if slide_data['slide_type'] == "title":
                # Create title slide
                create_title_slide(prs, slide_data['title'], "Professional Presentation")
                
            elif slide_data['slide_type'] == "section":
                # Create section divider
                create_section_slide(prs, slide_data['title'])
                section_count += 1
                
            else:
                # Create content slide with image
                create_content_slide_with_image(prs, slide_data)
                
                # Add section slides every 3-4 content slides
                if (i > 0 and i % 4 == 0 and section_count < 2):
                    section_title = f"Section {section_count + 1}"
                    create_section_slide(prs, section_title)
                    section_count += 1
        
        # Add a thank you slide at the end
        thank_you_slide = {
            'title': "Thank You",
            'content': ["Questions & Discussion", "Contact for more information"],
            'slide_type': "content",
            'image_query': "thank you business"
        }
        create_content_slide_with_image(prs, thank_you_slide)
        
        # Save presentation to a temporary file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            prs.save(temp_file.name)
            
            # Upload to Cloudinary
            cloudinary_url = upload_to_cloudinary(temp_file.name, presentation_id)
            
            if cloudinary_url:
                # Store URL in Redis
                store_presentation_url(presentation_id, cloudinary_url)
                
                # Clean up temporary file
                os.unlink(temp_file.name)
                
                return cloudinary_url
            else:
                raise Exception("Failed to upload presentation to Cloudinary")
        
    except Exception as e:
        logger.error(f"Error creating PowerPoint: {e}")
        raise Exception(f"Failed to create presentation: {str(e)}")

def generate_enhanced_presentation_content(topic: str, slide_count: int = 10) -> List[Dict]:
    """Generate enhanced presentation content with better structure and more content"""
    prompt = f"""Create a professional presentation outline for the topic: '{topic}'
    
    Requirements:
    - Create {slide_count} slides
    - I REPEAT: CREATE MINIMUM {slide_count} SLIDES
    - Include a title slide
    - Include 2-3 section divider slides
    - Include a thank you slide at the end
    - Each content slide should have:
      * A clear, concise title
      * 4-6 detailed bullet points
      * Each bullet point should be 1-2 sentences
      * Include relevant statistics or examples where appropriate
    - Make content engaging, informative, and professional
    - Structure the presentation with a clear flow:
      * Introduction
      * Main points
      * Supporting evidence
      * Conclusion
    - Include specific examples and data points
    - Use professional business language
    """
    
    try:
        slides = create_presentation_content(topic, slide_count=slide_count)
        
        # Enhance each slide with more content
        for slide in slides:
            if slide['slide_type'] == 'content':
                # Ensure minimum content length
                if len(slide['content']) < 4:
                    # Add more detailed points
                    additional_points = [
                        "Include relevant industry statistics",
                        "Provide specific examples and case studies",
                        "Discuss implementation strategies",
                        "Address potential challenges and solutions"
                    ]
                    slide['content'].extend(additional_points[:4 - len(slide['content'])])
                
                # Add image query based on title
                slide['image_query'] = slide['title']
        
        return slides
        
    except Exception as e:
        logger.error(f"Error generating presentation content: {e}")
        raise 